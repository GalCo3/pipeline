"""Demo producer: feeds every *-lexical service a batch of example Kafka messages.

Each source has an `examples/<service>.json` fixture holding ten message variants
that the service's model accepts — the legal shapes of that source's payload:
index / update / delete paths, source aliases vs model field names, optional
fields present, absent and explicitly null, and (for cargo) every file type the
text extractor handles plus the two failure paths.

Cargo examples also get their S3 object uploaded to MinIO first, since the
service reads the file the message points at.

Records are produced **without a Kafka key** — the real sources do not key their
records, so neither does the demo.
"""

from __future__ import annotations

import io
import json
import logging
import os
import re
import socket
import struct
import time
import zlib
from collections.abc import Callable
from datetime import UTC, datetime, timedelta
from pathlib import Path
from typing import Any

import boto3
from botocore.exceptions import ClientError
from confluent_kafka import Producer

logging.basicConfig(level=logging.INFO, format="%(asctime)s %(levelname)s %(message)s")
logger = logging.getLogger("demo-producer")

EXAMPLES_DIR = Path(__file__).parent / "examples"

BOOTSTRAP_SERVERS = os.environ.get("BOOTSTRAP_SERVERS", "kafka:9092")
S3_ENDPOINT = os.environ.get("S3_ENDPOINT", "http://minio:9000")
S3_ACCESS_KEY = os.environ.get("S3_ACCESS_KEY", "minioadmin")
S3_SECRET_KEY = os.environ.get("S3_SECRET_KEY", "minioadmin")
S3_BUCKET = os.environ.get("S3_BUCKET", "cargo-lexical")
# Comma-separated fixture names (`cargo-lexical,chat-messages-lexical`); empty = all of them.
SOURCES = [name for name in os.environ.get("SOURCES", "").split(",") if name.strip()]
# Every service parses timestamps with fromisoformat and falls back to dateutil,
# so this format and its Z-suffixed variant are both accepted.
DATETIME_FORMAT = "%Y-%m-%dT%H:%M:%S"
# 0 = produce the demo batch once and exit; >0 = keep re-producing every N seconds.
INTERVAL_SECONDS = float(os.environ.get("INTERVAL_SECONDS", "0"))
# Keeps document ids and S3 keys distinct across re-runs of the job, so a rerun
# adds documents instead of overwriting the previous batch.
RUN_SEED = int(os.environ.get("RUN_SEED", "") or zlib.crc32(socket.gethostname().encode()) % 90_000)

CARGO_SOURCES = {"cargo-operational-lexical", "cargo-my-storage-lexical"}


def build_txt() -> bytes:
    return (
        b"Cargo demo plain-text document.\n"
        b"Extracted by the streaming text extractor (no Tika needed).\n"
        b"Line three, with some searchable words: manifest, container, shipment.\n"
    )


def build_docx() -> bytes:
    from docx import Document

    document = Document()
    document.add_heading("Cargo Demo DOCX", level=1)
    document.add_paragraph("Shipment manifest for container HERMES-0001.")
    document.add_paragraph("Handled by the buffered DOCX extractor.")

    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def build_pdf() -> bytes:
    from reportlab.pdfgen import canvas

    buffer = io.BytesIO()
    pdf = canvas.Canvas(buffer)
    pdf.drawString(72, 800, "Cargo Demo PDF")
    pdf.drawString(72, 780, "Bill of lading for container HERMES-0002.")
    pdf.save()
    return buffer.getvalue()


def build_xlsx() -> bytes:
    from openpyxl import Workbook

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "manifest"
    sheet.append(["container", "weight_kg", "destination"])
    sheet.append(["HERMES-0003", 1200, "Haifa"])
    sheet.append(["HERMES-0004", 890, "Ashdod"])

    buffer = io.BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def build_pptx() -> bytes:
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Cargo Demo PPTX"
    slide.placeholders[1].text = "Port briefing for container HERMES-0005."

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def build_png() -> bytes:
    """A 1x1 PNG, assembled by hand so the image path needs no imaging library.

    Images are routed to the remote Tika extractor, which is what this example
    exercises — the pixel content is irrelevant.
    """

    def chunk(kind: bytes, payload: bytes) -> bytes:
        body = kind + payload
        return struct.pack(">I", len(payload)) + body + struct.pack(">I", zlib.crc32(body))

    header = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
    pixels = zlib.compress(b"\x00\xff\x00\x00")
    return (
        b"\x89PNG\r\n\x1a\n" + chunk(b"IHDR", header) + chunk(b"IDAT", pixels) + chunk(b"IEND", b"")
    )


def build_bin() -> bytes:
    """Binary noise: detected as application/octet-stream, which no extractor
    handles, so the service logs it and skips the message."""
    return bytes(range(256)) * 4


# Fixture `file` value -> (extension, content type, body builder).
FILE_BUILDERS: dict[str, tuple[str, str, Callable[[], bytes]]] = {
    "txt": ("txt", "text/plain", build_txt),
    "docx": (
        "docx",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        build_docx,
    ),
    "pdf": ("pdf", "application/pdf", build_pdf),
    "xlsx": (
        "xlsx",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        build_xlsx,
    ),
    "pptx": (
        "pptx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        build_pptx,
    ),
    "png": ("png", "image/png", build_png),
    "bin": ("bin", "application/octet-stream", build_bin),
}
# Fixture `file` value for the example whose object is deliberately absent.
MISSING_FILE = "missing"


def create_s3_client():
    return boto3.client(
        "s3",
        aws_access_key_id=S3_ACCESS_KEY,
        aws_secret_access_key=S3_SECRET_KEY,
        endpoint_url=S3_ENDPOINT,
    )


def ensure_bucket(s3) -> None:
    try:
        s3.head_bucket(Bucket=S3_BUCKET)
    except ClientError:
        s3.create_bucket(Bucket=S3_BUCKET)
        logger.info("Created bucket %s", S3_BUCKET)


def load_sources() -> list[tuple[str, str, list[dict]]]:
    """Returns (source name, topic, examples) per fixture, in file-name order."""
    names = SOURCES or sorted(path.stem for path in EXAMPLES_DIR.glob("*.json"))
    sources = []
    for name in names:
        fixture = json.loads((EXAMPLES_DIR / f"{name}.json").read_text(encoding="utf-8"))
        topic = os.environ.get(f"TOPIC_{name.upper().replace('-', '_')}", fixture["topic"])
        sources.append((name, topic, fixture["examples"]))
    return sources


def build_placeholders(doc_id: int, s3_key: str | None) -> dict[str, Any]:
    """Values the fixtures reference by `{{name}}`, resolved per example.

    A message must not mix naive and Z-suffixed timestamps: cargo and chief
    compare two of their own timestamps, and comparing naive to aware raises.
    """
    now = datetime.now(UTC)
    past = now - timedelta(days=7)
    future = now + timedelta(days=30)

    return {
        "id": str(doc_id),
        "id_int": doc_id,
        "now": now.strftime(DATETIME_FORMAT),
        "past": past.strftime(DATETIME_FORMAT),
        "future": future.strftime(DATETIME_FORMAT),
        "now_z": f"{now.strftime(DATETIME_FORMAT)}Z",
        "past_z": f"{past.strftime(DATETIME_FORMAT)}Z",
        "future_z": f"{future.strftime(DATETIME_FORMAT)}Z",
        "s3_key": s3_key,
        "s3_bucket": S3_BUCKET,
    }


# Only a whole string is substituted, so message bodies keep their own braces
# (and their `@mentions`, which an `@`-prefixed syntax would have swallowed).
PLACEHOLDER = re.compile(r"^\{\{(\w+)\}\}$")


def resolve(value: Any, placeholders: dict[str, Any]) -> Any:
    if isinstance(value, dict):
        return {key: resolve(item, placeholders) for key, item in value.items()}
    if isinstance(value, list):
        return [resolve(item, placeholders) for item in value]
    if isinstance(value, str) and (match := PLACEHOLDER.match(value)):
        name = match.group(1)
        if name not in placeholders:
            raise KeyError(f"Unknown placeholder {value} in example fixture")
        return placeholders[name]
    return value


def upload_example_file(s3, run_id: int, index: int, example: dict) -> str | None:
    """Uploads the example's S3 object under a per-run prefix. Returns its key."""
    file_kind = example.get("file")
    if file_kind is None:
        return None

    prefix = f"demo/run-{RUN_SEED}-{run_id}"
    if file_kind == MISSING_FILE:
        return f"{prefix}/{index:02d}-{example['name']}-does-not-exist.pdf"

    extension, content_type, build = FILE_BUILDERS[file_kind]
    key = f"{prefix}/{index:02d}-{example['name']}.{extension}"
    body = build()
    s3.put_object(Bucket=S3_BUCKET, Key=key, Body=body, ContentType=content_type)
    logger.info("Uploaded s3://%s/%s (%d bytes)", S3_BUCKET, key, len(body))
    return key


def delivery_report(err, msg) -> None:
    if err is not None:
        logger.error("Delivery to %s failed: %s", msg.topic(), err)
    else:
        logger.info("Produced to %s [%s] @ %s", msg.topic(), msg.partition(), msg.offset())


def build_ids(examples: list[dict], ordinal: int, run_id: int) -> list[int]:
    """One document id per example.

    An example may set `"id_of": "<other example name>"` to reuse that example's
    id — how the update and delete routes address a document an earlier example
    in the same batch just indexed, instead of an id Elasticsearch never saw.
    """
    ids = [
        RUN_SEED * 1_000_000 + run_id * 10_000 + ordinal * 100 + index
        for index in range(1, len(examples) + 1)
    ]
    by_name = {example["name"]: index for index, example in enumerate(examples)}

    for index, example in enumerate(examples):
        target = example.get("id_of")
        if target is not None:
            ids[index] = ids[by_name[target]]
    return ids


def produce_source(
    producer: Producer,
    source: str,
    topic: str,
    examples: list[dict],
    ordinal: int,
    run_id: int,
    s3,
) -> None:
    ids = build_ids(examples, ordinal, run_id)

    for index, example in enumerate(examples, start=1):
        doc_id = ids[index - 1]
        s3_key = upload_example_file(s3, run_id, index, example) if s3 else None
        message = resolve(example["message"], build_placeholders(doc_id, s3_key))

        # No key: the real sources produce keyless records.
        producer.produce(
            topic, value=json.dumps(message).encode("utf-8"), on_delivery=delivery_report
        )
        logger.info(
            "Queued %s example %d/%d %s (id=%s, expect=%s)",
            source,
            index,
            len(examples),
            example["name"],
            doc_id,
            example.get("expect", "indexed"),
        )


def main() -> None:
    sources = load_sources()
    logger.info("Producing %d sources: %s", len(sources), ", ".join(name for name, _, _ in sources))

    s3 = None
    if any(name in CARGO_SOURCES for name, _, _ in sources):
        s3 = create_s3_client()
        ensure_bucket(s3)

    producer = Producer({"bootstrap.servers": BOOTSTRAP_SERVERS})

    run_id = 0
    while True:
        run_id += 1
        for ordinal, (name, topic, examples) in enumerate(sources):
            produce_source(
                producer,
                name,
                topic,
                examples,
                ordinal,
                run_id,
                s3 if name in CARGO_SOURCES else None,
            )
            producer.flush(10)
            logger.info(
                "Batch %d: %d %s messages produced to %s", run_id, len(examples), name, topic
            )

        if INTERVAL_SECONDS <= 0:
            return
        time.sleep(INTERVAL_SECONDS)


if __name__ == "__main__":
    main()
